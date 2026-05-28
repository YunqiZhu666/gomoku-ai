"""
md2pptx.py — 将答辩PPT.md 转换为 PowerPoint 演示文稿
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色主题 ──
COLOR_PRIMARY = RGBColor(0x1A, 0x23, 0x7E)    # 深蓝
COLOR_ACCENT = RGBColor(0x0D, 0x47, 0xA1)     # 中蓝
COLOR_HIGHLIGHT = RGBColor(0xE8, 0x4D, 0x3D)  # 红色强调
COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)         # 浅灰蓝背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
COLOR_LIGHT_BLUE = RGBColor(0xEB, 0xF0, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def add_bg(slide, color=COLOR_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=COLOR_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=16, bold=False, color=COLOR_TEXT,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4),
                  font_name="Microsoft YaHei"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def add_accent_bar(slide, left, top, width, height, color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════
# Slide 1: 封面
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, COLOR_PRIMARY)
# 装饰条
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), COLOR_HIGHLIGHT)
add_shape_bg(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), COLOR_HIGHLIGHT)
# 主标题
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "五子棋 AI：从 Minimax 到 Alpha-Beta 剪枝",
            font_size=40, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
# 副标题
add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "人工智能基础 — 大作业答辩",
            font_size=22, color=RGBColor(0xA0, 0xC4, 0xFF), alignment=PP_ALIGN.CENTER)
# 分隔线
add_accent_bar(slide, Inches(5), Inches(4.0), Inches(3.333), Inches(0.04), COLOR_HIGHLIGHT)
# 署名
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
            "[姓名] · [学号] · 单人完成",
            font_size=18, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
# 底部信息
add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "2025 · 浙江大学",
            font_size=14, color=RGBColor(0x88, 0xAA, 0xDD), alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# Slide 2: 目录
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.8),
            "目  录", font_size=36, bold=True, color=COLOR_PRIMARY)

items = [
    ("01", "选题背景", "从深蓝到AlphaGo，经典博弈AI的智慧"),
    ("02", "核心原理", "Minimax → Alpha-Beta 剪枝"),
    ("03", "实现亮点", "评估函数 + 搜索优化 + Pygame界面"),
    ("04", "实验对比", "326×加速比，7层搜索深度"),
    ("05", "Demo 演示", "人机对战现场展示"),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.8 + i * 1.1
    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_ACCENT if i != 4 else COLOR_HIGHLIGHT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # 标题
    add_textbox(slide, Inches(2.2), Inches(y - 0.05), Inches(8), Inches(0.5),
                title, font_size=24, bold=True, color=COLOR_TEXT)
    # 描述
    add_textbox(slide, Inches(2.2), Inches(y + 0.4), Inches(8), Inches(0.4),
                desc, font_size=14, color=COLOR_GRAY)

# ══════════════════════════════════════════════
# Slide 3: 选题背景
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.6),
            "01  选题背景", font_size=28, bold=True, color=COLOR_PRIMARY)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), COLOR_HIGHLIGHT)

# 左侧quote
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(1),
            '"AI 不止有大模型，聪明地搜索本身就是一种强大的智能。"',
            font_size=18, bold=True, color=COLOR_ACCENT)

# 时间线
timeline = [
    ("1997", "深蓝 (IBM)\nMinimax + Alpha-Beta\n击败卡斯帕罗夫"),
    ("2016", "AlphaGo\nMCTS + 深度学习\n击败李世石"),
    ("2025", "本项目\n五子棋 AI\n9×9 棋盘"),
]
for i, (year, desc) in enumerate(timeline):
    x = 0.8 + i * 4.0
    add_shape_bg(slide, Inches(x), Inches(3.0), Inches(3.5), Inches(3.5), COLOR_WHITE)
    # 年份
    add_textbox(slide, Inches(x + 0.3), Inches(3.2), Inches(3), Inches(0.6),
                year, font_size=32, bold=True, color=COLOR_HIGHLIGHT)
    # 描述
    add_textbox(slide, Inches(x + 0.3), Inches(4.0), Inches(3), Inches(2.0),
                desc, font_size=14, color=COLOR_TEXT)

# 底部说明
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
            "选择五子棋：搜索复杂度适中 · 评估函数设计空间大 · 答辩演示效果好",
            font_size=14, color=COLOR_GRAY)

# ══════════════════════════════════════════════
# Slide 4: 核心原理
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
            "02  核心原理：Minimax → Alpha-Beta", font_size=28, bold=True, color=COLOR_PRIMARY)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), COLOR_HIGHLIGHT)

# Minimax 框
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
            "Minimax 算法", font_size=22, bold=True, color=COLOR_PRIMARY)
mm_items = [
    "▸ 我方 (MAX) 最大化分数，对手 (MIN) 最小化分数",
    "▸ 假设双方都走出最优棋",
    "▸ 递归搜索博弈树，深度 4 需搜索 19⁴ ≈ 130K 节点",
    "▸ 局限：指数级搜索爆炸，深度 4 需 15 秒",
]
for i, txt in enumerate(mm_items):
    add_textbox(slide, Inches(1.1), Inches(2.4 + i * 0.4), Inches(5), Inches(0.4),
                txt, font_size=13, color=COLOR_TEXT)

# Alpha-Beta 框
add_shape_bg(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), COLOR_WHITE)
add_textbox(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
            "Alpha-Beta 剪枝", font_size=22, bold=True, color=COLOR_HIGHLIGHT)
ab_items = [
    "▸ α = 我方能保证的最小分",
    "▸ β = 对方能接受的最大分",
    "▸ 当 α ≥ β，剪掉该分支",
    "▸ 效果：搜索节点从 b^d 降到 ≈ 2·b^(d/2)",
]
for i, txt in enumerate(ab_items):
    add_textbox(slide, Inches(7.3), Inches(2.4 + i * 0.4), Inches(5), Inches(0.4),
                txt, font_size=13, color=COLOR_TEXT)

# 评估函数框
add_shape_bg(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(4.6), Inches(5), Inches(0.5),
            "评估函数设计：棋型识别", font_size=22, bold=True, color=COLOR_PRIMARY)
# 权重表
weights = [
    ("连五", "10,000,000", COLOR_HIGHLIGHT),
    ("活四", "8,000,000", COLOR_ORANGE),
    ("冲四 / 活三", "2,500,000", COLOR_ACCENT),
    ("眠三", "200,000", COLOR_GREEN),
    ("活二", "5,000", COLOR_GRAY),
]
for i, (name, score, color) in enumerate(weights):
    x = 1.1 + i * 2.3
    box = add_shape_bg(slide, Inches(x), Inches(5.3), Inches(2.0), Inches(1.0), color)
    add_textbox(slide, Inches(x + 0.1), Inches(5.35), Inches(1.8), Inches(0.4),
                name, font_size=14, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(x + 0.1), Inches(5.7), Inches(1.8), Inches(0.4),
                score, font_size=16, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.1), Inches(6.5), Inches(11), Inches(0.4),
            "最终分数 = 己方棋型总分 − 对方棋型总分 × 1.1（防守偏重）",
            font_size=13, color=COLOR_GRAY)

# ══════════════════════════════════════════════
# Slide 5: 实现亮点
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
            "03  实现亮点", font_size=28, bold=True, color=COLOR_PRIMARY)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), COLOR_HIGHLIGHT)

# 技术栈
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(2.2), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(1.6), Inches(3), Inches(0.4),
            "技术栈", font_size=20, bold=True, color=COLOR_PRIMARY)
add_textbox(slide, Inches(1.1), Inches(2.1), Inches(3), Inches(1.2),
            "Python 3  ·  Pygame  ·  NumPy\nMatplotlib  ·  python-pptx\n\n9×9 棋盘 · ~2000 行代码",
            font_size=14, color=COLOR_TEXT)

# 架构
add_shape_bg(slide, Inches(4.7), Inches(1.5), Inches(3.5), Inches(2.2), COLOR_WHITE)
add_textbox(slide, Inches(5.0), Inches(1.6), Inches(3), Inches(0.4),
            "模块架构", font_size=20, bold=True, color=COLOR_PRIMARY)
modules = [
    ("board.py", "棋盘引擎"),
    ("evaluate.py", "评估函数"),
    ("search.py", "搜索算法"),
    ("ui.py", "图形界面"),
]
for i, (f, desc) in enumerate(modules):
    add_textbox(slide, Inches(5.0), Inches(2.2 + i * 0.35), Inches(3), Inches(0.35),
                f"📄 {f} — {desc}", font_size=12, color=COLOR_TEXT)

# 优化技术
add_shape_bg(slide, Inches(8.6), Inches(1.5), Inches(4.0), Inches(2.2), COLOR_WHITE)
add_textbox(slide, Inches(8.9), Inches(1.6), Inches(3.5), Inches(0.4),
            "三大优化", font_size=20, bold=True, color=COLOR_PRIMARY)
opts = [
    ("走法裁剪", "周围 1 格，81 → ~15"),
    ("走法排序", "剪枝效率提升 2-3×"),
    ("置换表", "缓存已搜局面"),
]
for i, (opt, desc) in enumerate(opts):
    add_textbox(slide, Inches(8.9), Inches(2.2 + i * 0.4), Inches(3.5), Inches(0.35),
                f"✅ {opt}：{desc}", font_size=13, color=COLOR_TEXT)

# Bug修复亮点
bug_box = add_shape_bg(slide, Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.7), COLOR_PRIMARY)
add_textbox(slide, Inches(1.1), Inches(3.95), Inches(11.3), Inches(0.55),
            "🐛 评估函数调试：修复了三个关键 Bug（滑动窗口重复计数、走法排序视角错配、跳连分组对手棋子误判），修复后 AI 在任意搜索深度均能正确封堵活三",
            font_size=11, color=COLOR_WHITE)

# 下部分：文件结构和关键功能
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.5), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(0.5),
            "项目文件结构", font_size=20, bold=True, color=COLOR_PRIMARY)
add_textbox(slide, Inches(1.1), Inches(5.4), Inches(11), Inches(1.8),
            "gomoku-ai/\n"
            "├── board.py              棋盘逻辑（落子 / 胜负 / 走法生成）\n"
            "├── evaluate.py           棋型识别 + 权重打分\n"
            "├── search.py             Minimax → Alpha-Beta → 迭代加深\n"
            "├── ui.py                 Pygame 图形对战界面\n"
            "├── main.py               统一入口 (GUI/CLI/实验)\n"
            "├── experiment.py / experiment_defense.py   实验对比框架\n"
            "├── assets/               实验图表\n"
            "├── 开题报告.md / 实验报告.md / 答辩PPT.md",
            font_size=12, color=COLOR_TEXT)

# ══════════════════════════════════════════════
# Slide 6: 实验对比
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
            "04  实验对比", font_size=28, bold=True, color=COLOR_PRIMARY)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), COLOR_HIGHLIGHT)

# 表格数据
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(7.0), Inches(3.5), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(1.6), Inches(6.5), Inches(0.4),
            "Minimax vs Alpha-Beta 耗时对比", font_size=18, bold=True, color=COLOR_PRIMARY)

# 手动画表格
table_data = [
    ["搜索深度", "Minimax", "Alpha-Beta", "加速比"],
    ["1", "0.001s", "0.001s", "1×"],
    ["2", "0.020s", "0.005s", "4×"],
    ["3", "0.500s", "0.016s", "31×"],
    ["4", "15.000s", "0.046s", "326×"],
]
row_h = Inches(0.45)
col_w = [Inches(1.2), Inches(1.6), Inches(1.8), Inches(1.6)]
start_x = Inches(1.1)
start_y = Inches(2.2)
for ri, row in enumerate(table_data):
    for ci, cell in enumerate(row):
        x = start_x + sum(col_w[:ci])
        y = start_y + ri * row_h
        box = add_shape_bg(slide, x, y, col_w[ci], row_h,
                           COLOR_LIGHT_BLUE if ri == 0 else COLOR_WHITE)
        fc = COLOR_WHITE if ri == 0 else COLOR_TEXT
        add_textbox(slide, x + Inches(0.05), y + Inches(0.05), col_w[ci] - Inches(0.1), row_h - Inches(0.1),
                    cell, font_size=13 if ri > 0 else 14,
                    bold=(ri == 0 or ci == 3),
                    color=fc, alignment=PP_ALIGN.CENTER)

# 右侧关键数字
add_shape_bg(slide, Inches(8.2), Inches(1.5), Inches(4.3), Inches(3.5), COLOR_PRIMARY)
add_textbox(slide, Inches(8.5), Inches(1.7), Inches(3.7), Inches(0.5),
            "关键指标", font_size=20, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
kpis = [
    ("326×", "深度 4 加速比"),
    ("7 层", "3 秒内搜索深度"),
    ("70%", "先手胜率"),
    ("~15", "平均分支因子"),
]
for i, (num, desc) in enumerate(kpis):
    y = 2.4 + i * 0.75
    add_textbox(slide, Inches(8.5), Inches(y), Inches(1.5), Inches(0.5),
                num, font_size=28, bold=True, color=COLOR_HIGHLIGHT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.0), Inches(y + 0.1), Inches(2.5), Inches(0.4),
                desc, font_size=13, color=RGBColor(0xCC, 0xDD, 0xFF))

# 底部：防守权重实验（紧凑版）
add_shape_bg(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.8), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(5.4), Inches(11), Inches(0.4),
            "防守权重对弈实验：不同防守系数 AI 对战胜率（vs baseline 1.0）", font_size=16, bold=True, color=COLOR_PRIMARY)
defense_data = [
    ["防守权重", "0.5", "0.8", "1.0", "1.5", "2.0", "3.0"],
    ["胜率", "33%", "50%", "50%", "50%", "67%", "83%"],
]
for ri, row in enumerate(defense_data):
    for ci, cell in enumerate(row):
        x = Inches(1.1) + ci * Inches(1.65)
        y = Inches(5.95) + ri * Inches(0.5)
        box = add_shape_bg(slide, x, y, Inches(1.5), Inches(0.45),
                           COLOR_LIGHT_BLUE if ri == 0 else COLOR_WHITE)
        fc = COLOR_WHITE if ri == 0 else (COLOR_HIGHLIGHT if ci == 6 else COLOR_TEXT)
        add_textbox(slide, x + Inches(0.05), y + Inches(0.05), Inches(1.4), Inches(0.35),
                    cell, font_size=13,
                    bold=(ri == 0 or ci == 6),
                    color=fc, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.1), Inches(6.78), Inches(11), Inches(0.3),
            "→ 防守权重 3.0 胜率最高 (83%)，验证了加强防守能有效提升棋力；权重 1.1 兼顾攻守",
            font_size=11, color=COLOR_GRAY)

# ══════════════════════════════════════════════
# Slide 7: 搜索树对比 + 评估函数灵敏度
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
            "04  实验对比（续）：搜索树 & 棋型灵敏度", font_size=28, bold=True, color=COLOR_PRIMARY)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), COLOR_HIGHLIGHT)

# 左侧：搜索树大小对比
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(3.0), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(1.6), Inches(5.5), Inches(0.4),
            "搜索树大小对比", font_size=20, bold=True, color=COLOR_PRIMARY)
tree_data = [
    ["深度 4", "130,321 节点", "400 节点", "326×"],
    ["深度 5", "2,476,099 节点", "900 节点", "2,751×"],
    ["深度 6", "47,045,881 节点", "1,800 节点", "26,137×"],
]
for ri, row in enumerate(tree_data):
    for ci, cell in enumerate(row):
        x = Inches(1.1) + ci * Inches(1.4)
        y = Inches(2.2) + ri * Inches(0.5)
        c = COLOR_HIGHLIGHT if ci == 3 else COLOR_TEXT
        add_textbox(slide, x, y, Inches(1.3), Inches(0.4),
                    cell, font_size=13, bold=(ci == 3), color=c)
# 表头
for ci, h in enumerate(["深度", "Minimax", "Alpha-Beta", "缩减"]):
    add_textbox(slide, Inches(1.1) + ci * Inches(1.4), Inches(2.2), Inches(1.3), Inches(0.3),
                h, font_size=11, bold=True, color=COLOR_GRAY)
add_textbox(slide, Inches(1.1), Inches(3.8), Inches(5.5), Inches(0.4),
            "深度 6 时 Minimax 需 4700 万节点，Alpha-Beta 仅需 1800 节点 → 可多搜 3-4 层",
            font_size=12, color=COLOR_GRAY)

# 右侧：评估函数灵敏度
add_shape_bg(slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(3.0), COLOR_WHITE)
add_textbox(slide, Inches(7.5), Inches(1.6), Inches(4.8), Inches(0.4),
            "评估函数棋型灵敏度", font_size=20, bold=True, color=COLOR_PRIMARY)
sensitivity = [
    ("空棋盘", "0", "基准"),
    ("一子", "40", "微弱先手"),
    ("活二", "1,060", "潜在威胁"),
    ("活三", "100,090", "威胁大"),
    ("活四", "1,000,120", "接近必胜"),
    ("连五", "10,000,150", "必胜"),
]
for i, (name, score, meaning) in enumerate(sensitivity):
    y = Inches(2.2 + i * 0.42)
    add_textbox(slide, Inches(7.5), y, Inches(1.4), Inches(0.35),
                name, font_size=13, bold=True, color=COLOR_TEXT)
    add_textbox(slide, Inches(8.9), y, Inches(1.6), Inches(0.35),
                score, font_size=12, color=COLOR_ACCENT)
    add_textbox(slide, Inches(10.5), y, Inches(2.0), Inches(0.35),
                meaning, font_size=11, color=COLOR_GRAY)

# 底部：AI自对弈
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(4.9), Inches(11), Inches(0.4),
            "AI vs AI 自对弈 & 综合结论", font_size=20, bold=True, color=COLOR_PRIMARY)
selfplay = [
    "▸ 先手（黑）胜率 70%，后手（白）胜率 20%，平局 10%，平均 23.5 手",
    "▸ Alpha-Beta 剪枝效果显著：深度 ≥ 3 时加速比 30× 以上，相同时间可搜更深 3-4 层",
    "▸ 防守权重实验：权重 3.0 胜率 83%，验证加强防守有效提升棋力；权重 1.1 兼顾攻守",
    "▸ 走法排序 + 置换表 + 迭代加深 + 走法裁剪的组合效果远超单一技术",
]
for i, txt in enumerate(selfplay):
    add_textbox(slide, Inches(1.1), Inches(5.5 + i * 0.35), Inches(11), Inches(0.35),
                txt, font_size=13, color=COLOR_TEXT)

# ══════════════════════════════════════════════
# Slide 8: Demo + 总结
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, COLOR_ACCENT)
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
            "05  Demo 演示 & 总结", font_size=28, bold=True, color=COLOR_PRIMARY)
add_accent_bar(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.04), COLOR_HIGHLIGHT)

# 左侧：Demo
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.0), COLOR_WHITE)
add_textbox(slide, Inches(1.1), Inches(1.6), Inches(5.3), Inches(0.5),
            "🎮 Demo 演示", font_size=22, bold=True, color=COLOR_PRIMARY)
demo_items = [
    "1. 运行: python main.py",
    "2. 玩家执黑先行，鼠标点击落子",
    "3. AI 自动思考并应对",
    "4. 右侧信息栏显示搜索统计",
    "",
    "展示重点:",
    "✅ AI 识别活三 → 正确堵截",
    "✅ AI 识别冲四 → 优先防守",
    "✅ 一步杀 → 立即取胜",
    "✅ 悔棋功能",
]
for i, txt in enumerate(demo_items):
    c = COLOR_HIGHLIGHT if i >= 6 else COLOR_TEXT
    add_textbox(slide, Inches(1.1), Inches(2.3 + i * 0.38), Inches(5.3), Inches(0.35),
                txt, font_size=14, bold=(i >= 6), color=c)

# 右侧：总结
add_shape_bg(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0), COLOR_PRIMARY)
add_textbox(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
            "📌 总结", font_size=22, bold=True, color=COLOR_WHITE)

summary_items = [
    "实现了完整五子棋 AI：",
    "  ✓ Minimax + Alpha-Beta 剪枝",
    "  ✓ 迭代加深搜索（7 层/3秒）",
    "  ✓ Pygame 图形对战界面",
    "  ✓ 置换表 + 走法排序优化",
    "",
    "实验验证：",
    "  ✓ Alpha-Beta 加速比 326×",
    "  ✓ 防守权重 3.0 胜率 83%",
    "  ✓ 评估函数准确识别棋型",
    "  ✓ 修复 3 个关键 Bug",
    "",
    "理解了博弈树搜索的经典范式",
    "与算法优化和调试的工程实践",
]
for i, txt in enumerate(summary_items):
    c = COLOR_WHITE
    add_textbox(slide, Inches(7.3), Inches(2.4 + i * 0.35), Inches(5), Inches(0.35),
                txt, font_size=13, bold=("✓" in txt or "：" in txt), color=c)

# ══════════════════════════════════════════════
# Slide 9: 谢谢
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_PRIMARY)
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), COLOR_HIGHLIGHT)
add_shape_bg(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), COLOR_HIGHLIGHT)
add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
            "谢谢！", font_size=60, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(5.5), Inches(3.3), Inches(2.333), Inches(0.04), COLOR_HIGHLIGHT)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
            "[姓名] · [学号]", font_size=22, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
            "Q & A", font_size=28, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = "答辩PPT.pptx"
prs.save(output_path)
print(f"[OK] PPT 已生成: {output_path}")
print(f"     共 {len(prs.slides)} 页幻灯片")
